from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Diligent GTM Engineer Analysis Project"
subtitle.text = "Building a Data-Driven ICP Framework and Account Prioritization System\nDate: September 15, 2025"

# Slide 2: Project Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Project Overview"
tf = body_shape.text_frame
tf.text = "Objectives:"
p = tf.add_paragraph()
p.text = "• Data Processing & Cleaning"
p.level = 1
p = tf.add_paragraph()
p.text = "• Pipeline/Script Development"
p.level = 1
p = tf.add_paragraph()
p.text = "• ICP Definition & Prioritization"
p.level = 1
p = tf.add_paragraph()
p.text = "• Systems Integration Plan"
p.level = 1

# Slide 3: Data Cleaning Approach
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Data Cleaning Approach"
tf = body_shape.text_frame
tf.text = "Key Cleaning Tasks:"
p = tf.add_paragraph()
p.text = "• Employee Count Normalization: Standardized ranges and formats"
p.level = 1
p = tf.add_paragraph()
p.text = "• Revenue Standardization: Universal USD format"
p.level = 1
p = tf.add_paragraph()
p.text = "• Region Mapping: Consistent labels and expansions"
p.level = 1
p = tf.add_paragraph()
p.text = "• Date Normalization: YYYY-MM-DD format"
p.level = 1

# Slide 4: ICP Scoring Framework
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "ICP Scoring Framework"
tf = body_shape.text_frame
tf.text = "ICP Archetypes:"
p = tf.add_paragraph()
p.text = "• Enterprise Risk Management: Large companies focused on risk"
p.level = 1
p = tf.add_paragraph()
p.text = "• Mid-Market Compliance: Growing companies needing compliance"
p.level = 1
p = tf.add_paragraph()
p.text = "• Board Governance: Organizations managing board processes"
p.level = 1
p = tf.add_paragraph()
p.text = "Scoring: 0-100 points across firmographic, solution, intent, and tech dimensions"
p.level = 0

# Slide 5: Key Results and Next Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Key Results and Next Steps"
tf = body_shape.text_frame
tf.text = "Results:"
p = tf.add_paragraph()
p.text = "• 20 Critical Accounts (90+ points) for immediate outreach"
p.level = 1
p = tf.add_paragraph()
p.text = "• Data quality improved from 80% to 95%+ completeness"
p.level = 1
p = tf.add_paragraph()
p.text = "• Clear prioritization framework for 1,000 prospects"
p.level = 1
p = tf.add_paragraph()
p.text = "Next Steps:"
p.level = 0
p = tf.add_paragraph()
p.text = "• Import scored data into Salesforce"
p.level = 1
p = tf.add_paragraph()
p.text = "• Begin outreach to top 20 critical accounts"
p.level = 1
p = tf.add_paragraph()
p.text = "• Build automated scoring pipeline"
p.level = 1

# Save the presentation
prs.save('/Users/saipavankatineedi/Desktop/Diligent_project/deliverables/project_presentation.pptx')
